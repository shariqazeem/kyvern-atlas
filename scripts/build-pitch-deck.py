#!/usr/bin/env python3
"""
build-pitch-deck.py — generates Kyvern's Demo Day pitch deck.

7 custom-designed slides matching the /app aesthetic:
  - Cream backdrop (#FAFAFA)
  - Charcoal text (#0A0A0A) / muted gray (#6B6B6B)
  - JetBrains Mono for every number, address, code
  - Inter for body
  - Hairline dividers, generous negative space, "device" feel

Run:
  python3 scripts/build-pitch-deck.py

Output:
  ./kyvern-pitch-deck.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ────────────────────────────────────────────────────────────────────────
# Design tokens — match /app exactly
# ────────────────────────────────────────────────────────────────────────

CREAM = RGBColor(0xFA, 0xFA, 0xFA)
INK = RGBColor(0x0A, 0x0A, 0x0A)
INK_2 = RGBColor(0x1F, 0x29, 0x37)
INK_3 = RGBColor(0x6B, 0x6B, 0x6B)
INK_4 = RGBColor(0xA0, 0xA0, 0xA0)
HAIRLINE = RGBColor(0xE5, 0xE5, 0xE5)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
RED = RGBColor(0xDC, 0x26, 0x26)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_SANS = "Inter"  # fallback Helvetica / system default
FONT_MONO = "JetBrains Mono"  # fallback Menlo (mac) / Consolas (win)
FONT_DISPLAY = "Bricolage Grotesque"  # used on landing — premium feel

REPO = Path(__file__).resolve().parents[1]
ASSETS = REPO / "assets" / "deck"
DECKS = REPO / "decks"
DECKS.mkdir(exist_ok=True)
OUT = DECKS / "kyvern-pitch-deck.pptx"

# 16:9 in inches
SLIDE_W = 13.333
SLIDE_H = 7.5

# ────────────────────────────────────────────────────────────────────────
# Helpers
# ────────────────────────────────────────────────────────────────────────


def set_solid_bg(slide, color):
    """Fill the slide background with a solid color."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H)
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    # send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(slide, x, y, w, h, text, *, font=FONT_SANS, size=14, color=INK,
             bold=False, align=PP_ALIGN.LEFT, line_spacing=1.2, italic=False,
             tracking=None):
    """Add a text box with one paragraph."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size)
    f.color.rgb = color
    f.bold = bold
    f.italic = italic
    return tb


def add_rule(slide, x, y, w, color=HAIRLINE, thickness_pt=0.75):
    """Add a horizontal hairline rule."""
    line = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y))
    line.line.color.rgb = color
    line.line.width = Pt(thickness_pt)
    return line


def add_pill(slide, x, y, w, h, text, *, fg=INK, bg=WHITE, border=HAIRLINE,
             font=FONT_MONO, size=9, bold=False):
    """Rounded-rect pill with text."""
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    pill.adjustments[0] = 0.5  # max-rounded
    pill.fill.solid()
    pill.fill.fore_color.rgb = bg
    pill.line.color.rgb = border
    pill.line.width = Pt(0.5)
    tf = pill.text_frame
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = fg
    r.font.bold = bold
    return pill


def add_card(slide, x, y, w, h, *, bg=WHITE, border=HAIRLINE,
             radius_pct=0.05, shadow=False):
    """White rounded card (like the /app cards)."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    card.adjustments[0] = radius_pct
    card.fill.solid()
    card.fill.fore_color.rgb = bg
    card.line.color.rgb = border
    card.line.width = Pt(0.5)
    card.shadow.inherit = False
    return card


def add_slide_marker(slide, number_str, label):
    """Tiny top-right slide ID + bottom-left slide label."""
    # top-right slide counter "01 / 07"
    add_text(slide, SLIDE_W - 1.5, 0.32, 1.4, 0.25,
             number_str, font=FONT_MONO, size=8, color=INK_4,
             align=PP_ALIGN.RIGHT, tracking=0.15)
    # bottom-left label "SLIDE 01 · COVER"
    add_text(slide, 0.5, SLIDE_H - 0.45, 6, 0.25,
             f"SLIDE {number_str.split('/')[0].strip()} · {label.upper()}",
             font=FONT_MONO, size=8, color=INK_4, align=PP_ALIGN.LEFT)


def add_top_brand_line(slide):
    """Top edge: 'BUILD ON SOLANA WITH KAST · PAKISTAN' right, tiny logo left."""
    # tiny KV mark left
    logo = ASSETS / "kyvernlabs_logo.jpg"
    if logo.exists():
        slide.shapes.add_picture(str(logo), Inches(0.5), Inches(0.30),
                                 height=Inches(0.30), width=Inches(0.30))
    # Track top-right
    add_text(slide, SLIDE_W - 6.5, 0.34, 5, 0.25,
             "BUILD ON SOLANA WITH KAST · PAKISTAN",
             font=FONT_MONO, size=8, color=INK_4, align=PP_ALIGN.RIGHT)


def add_bottom_url_line(slide, left="kyvernlabs.com  ·  @shariqshkt",
                        right="Shariq Shaukat"):
    add_text(slide, 0.5, SLIDE_H - 0.45, 6, 0.25, left,
             font=FONT_MONO, size=8, color=INK_4, align=PP_ALIGN.LEFT)
    add_text(slide, SLIDE_W - 6.5, SLIDE_H - 0.45, 6, 0.25, right,
             font=FONT_MONO, size=8, color=INK_4, align=PP_ALIGN.RIGHT)


# ────────────────────────────────────────────────────────────────────────
# Build the deck
# ────────────────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)

BLANK = prs.slide_layouts[6]


# ════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER (clean, single message)
# ════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
set_solid_bg(s1, CREAM)
add_top_brand_line(s1)
add_rule(s1, 0.5, 0.76, SLIDE_W - 1.0)

# Vertical-center the wordmark + tagline (lots of negative space)
# 96pt is enough to dominate without colliding with the line below.
add_text(s1, 0.5, 2.0, 12, 1.6, "Kyvern",
         font=FONT_DISPLAY, size=96, color=INK, bold=True,
         line_spacing=1.0)

# Single tagline below — the only message on the cover
add_text(s1, 0.5, 3.55, 12.3, 0.7,
         "Solana's policy layer for AI agents.",
         font=FONT_DISPLAY, size=28, color=INK_2, italic=True)

# Mono eyebrow under tagline
add_text(s1, 0.5, 4.25, 12.3, 0.4,
         "PER-AGENT WALLETS · ENFORCED ON-CHAIN BEFORE A DOLLAR MOVES",
         font=FONT_MONO, size=11, color=INK_4)

# Live numbers strip — proof, not promises
add_rule(s1, 0.5, 5.85, SLIDE_W - 1.0)
add_text(s1, 0.5, 6.00, 12, 0.3,
         "LIVE ON SOLANA DEVNET",
         font=FONT_MONO, size=9, color=INK_4, bold=True)

# 4 column number row
col_y = 6.30
labels = [
    ("39d", "live uptime"),
    ("17,339", "agent attempts"),
    ("6,905", "blocked on-chain"),
    ("$0.00", "lost"),
]
col_w = (SLIDE_W - 1.0) / 4
for i, (n, lbl) in enumerate(labels):
    cx = 0.5 + i * col_w
    add_text(s1, cx, col_y, col_w, 0.5, n,
             font=FONT_MONO, size=28, color=INK, bold=True, align=PP_ALIGN.LEFT)
    add_text(s1, cx, col_y + 0.62, col_w, 0.3, lbl,
             font=FONT_MONO, size=9, color=INK_4, align=PP_ALIGN.LEFT)

add_bottom_url_line(s1)
add_text(s1, SLIDE_W - 1.5, 0.32, 1.4, 0.25, "01 / 07",
         font=FONT_MONO, size=8, color=INK_4, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE HOOK (problem)
# ════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
set_solid_bg(s2, CREAM)
add_top_brand_line(s2)
add_rule(s2, 0.5, 0.76, SLIDE_W - 1.0)

# Eyebrow
add_text(s2, 0.5, 1.3, 8, 0.4,
         "SLIDE 02 · THE PROBLEM",
         font=FONT_MONO, size=10, color=INK_4)

# Big problem statement
add_text(s2, 0.5, 2.0, 12.3, 2.5,
         "Every AI agent that touches\nmoney today uses\nsomebody's keys.",
         font=FONT_DISPLAY, size=58, color=INK, bold=True, line_spacing=1.05)

# Detail line
add_text(s2, 0.5, 5.2, 12.3, 0.6,
         "One bad inference. One prompt injection. One runaway loop.\n"
         "Your wallet is empty.",
         font=FONT_SANS, size=19, color=INK_3, line_spacing=1.4)

add_bottom_url_line(s2)
add_text(s2, SLIDE_W - 1.5, 0.32, 1.4, 0.25, "02 / 07",
         font=FONT_MONO, size=8, color=INK_4, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION (in one frame)
# ════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
set_solid_bg(s3, CREAM)
add_top_brand_line(s3)
add_rule(s3, 0.5, 0.76, SLIDE_W - 1.0)

# Eyebrow
add_text(s3, 0.5, 1.3, 8, 0.4,
         "SLIDE 03 · THE SOLUTION",
         font=FONT_MONO, size=10, color=INK_4)

# Big payoff
add_text(s3, 0.5, 2.0, 12.3, 2.0,
         "Kyvern gives your agent\na Visa with a daily cap.",
         font=FONT_DISPLAY, size=56, color=INK, bold=True, line_spacing=1.05)

# Three short pillar lines
pillars = [
    ("BUDGETS", "Per-tx, daily, weekly USDC caps."),
    ("ALLOWLISTS", "Pre-approved merchants only."),
    ("KILL SWITCH", "One tap. Agent stops. On-chain."),
]
py = 4.65
pw = (SLIDE_W - 1.0) / 3
for i, (head, sub) in enumerate(pillars):
    px = 0.5 + i * pw
    add_text(s3, px, py, pw - 0.3, 0.4, head,
             font=FONT_MONO, size=11, color=INK, bold=True)
    add_text(s3, px, py + 0.45, pw - 0.3, 0.6, sub,
             font=FONT_SANS, size=15, color=INK_3, line_spacing=1.3)

# Bottom enforcement line
add_rule(s3, 0.5, 5.95, SLIDE_W - 1.0)
add_text(s3, 0.5, 6.10, 12.3, 0.4,
         "Enforced on Solana by two programs composed atomically — "
         "Kyvern policy + Squads v4. The chain is the arbiter. Not our server.",
         font=FONT_SANS, size=12, color=INK_3, italic=True, line_spacing=1.4)

add_bottom_url_line(s3)
add_text(s3, SLIDE_W - 1.5, 0.32, 1.4, 0.25, "03 / 07",
         font=FONT_MONO, size=8, color=INK_4, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 4 — LIVE DEMO PLACEHOLDER
# ════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
set_solid_bg(s4, CREAM)
add_top_brand_line(s4)
add_rule(s4, 0.5, 0.76, SLIDE_W - 1.0)

# Center the demo marker
# Eyebrow
add_text(s4, 0.5, 2.6, 12.3, 0.4,
         "SLIDE 04 · LIVE",
         font=FONT_MONO, size=11, color=INK_4, align=PP_ALIGN.CENTER)

# Big arrow + label
add_text(s4, 0.5, 3.05, 12.3, 1.6,
         "→ kyvernlabs.com / app",
         font=FONT_DISPLAY, size=62, color=INK, bold=True,
         align=PP_ALIGN.CENTER)

add_text(s4, 0.5, 4.6, 12.3, 0.5,
         "Let me show you instead of telling you.",
         font=FONT_SANS, size=20, color=INK_3, italic=True,
         align=PP_ALIGN.CENTER)

# bottom mini hints (what we'll click)
hints = [
    "01 · Run prediction agent (ParallaxPay)",
    "02 · Try $5 over-cap",
    "03 · Open Developer mode (live SDK events)",
]
hy = 5.7
for i, h in enumerate(hints):
    add_text(s4, 0.5, hy + i * 0.32, 12.3, 0.3, h,
             font=FONT_MONO, size=10, color=INK_4, align=PP_ALIGN.CENTER)

add_bottom_url_line(s4)
add_text(s4, SLIDE_W - 1.5, 0.32, 1.4, 0.25, "04 / 07",
         font=FONT_MONO, size=8, color=INK_4, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 5 — HOW IT WORKS (3 boxes + bottom proof line)
# ════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
set_solid_bg(s5, CREAM)
add_top_brand_line(s5)
add_rule(s5, 0.5, 0.76, SLIDE_W - 1.0)

add_text(s5, 0.5, 1.3, 8, 0.4, "SLIDE 05 · HOW IT WORKS",
         font=FONT_MONO, size=10, color=INK_4)

add_text(s5, 0.5, 1.85, 12.3, 0.7,
         "Three pieces. Four lines of code. Four hundred milliseconds.",
         font=FONT_DISPLAY, size=30, color=INK, bold=True, line_spacing=1.1)

# Three vertical cards
cx_left = 0.5
total_w = SLIDE_W - 1.0
card_gap = 0.3
card_w = (total_w - 2 * card_gap) / 3
card_y = 3.2
card_h = 3.0

for i, (head, code, sub) in enumerate([
    ("01 · YOUR AGENT",
     "vault.pay({\n  merchant,\n  amount,\n  memo,\n});",
     "Four lines of SDK. Drop into LangChain, Eliza, Claude Agent SDK."),
    ("02 · KYVERN ON-CHAIN",
     "PpmZ…MSqc\nrules.check(\n  budget · allowlist\n  velocity · memo\n);",
     "Refuses on-chain. Real failure code. Real signature. Real Explorer link."),
    ("03 · SQUADS SETTLES",
     "SQDS4ep…pCf\nspending_limit_use\n  → moves USDC\n  (or reverts)",
     "Audited by Trail of Bits, OtterSec, Neodyme. Secures $10B+ on Solana."),
]):
    cx = cx_left + i * (card_w + card_gap)
    add_card(s5, cx, card_y, card_w, card_h)
    # Card header eyebrow
    add_text(s5, cx + 0.25, card_y + 0.25, card_w - 0.5, 0.3,
             head, font=FONT_MONO, size=10, color=INK_4, bold=True)
    # Code block (mono) — tighter line spacing, smaller font, bounded box
    add_text(s5, cx + 0.25, card_y + 0.75, card_w - 0.5, 1.45,
             code, font=FONT_MONO, size=12, color=INK, line_spacing=1.35)
    # Caption (sans) — moved further down to clear the code block
    add_text(s5, cx + 0.25, card_y + 2.30, card_w - 0.5, 0.7,
             sub, font=FONT_SANS, size=10.5, color=INK_3, line_spacing=1.4)

# Bottom proof bar
add_rule(s5, 0.5, 6.40, SLIDE_W - 1.0)
add_text(s5, 0.5, 6.55, 12.3, 0.4,
         "Atomic CPI. Either rejects → entire tx reverts. No partial state. "
         "No off-chain trust.",
         font=FONT_SANS, size=12, color=INK_3, italic=True, line_spacing=1.4)

add_bottom_url_line(s5)
add_text(s5, SLIDE_W - 1.5, 0.32, 1.4, 0.25, "05 / 07",
         font=FONT_MONO, size=8, color=INK_4, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 6 — WHAT'S LIVE + WHAT'S NEXT
# ════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
set_solid_bg(s6, CREAM)
add_top_brand_line(s6)
add_rule(s6, 0.5, 0.76, SLIDE_W - 1.0)

add_text(s6, 0.5, 1.3, 8, 0.4, "SLIDE 06 · LIVE + ROADMAP",
         font=FONT_MONO, size=10, color=INK_4)

add_text(s6, 0.5, 1.85, 12.3, 0.7,
         "Shipped. Not promised.",
         font=FONT_DISPLAY, size=32, color=INK, bold=True, line_spacing=1.0)

# Two-column layout
left_x = 0.5
right_x = 7.0
col_w_half = 5.83  # leaves the gutter

# LEFT — What's live (numbers + bullets)
add_text(s6, left_x, 2.85, col_w_half, 0.3,
         "WHAT'S LIVE TODAY",
         font=FONT_MONO, size=10, color=INK_4, bold=True)

# 2x2 numbers
nums = [
    ("39d", "Atlas uptime"),
    ("17,339", "agent attempts"),
    ("6,905", "blocked on-chain"),
    ("$0.00", "lost"),
]
for i, (n, lbl) in enumerate(nums):
    row = i // 2
    col = i % 2
    nx = left_x + col * (col_w_half / 2)
    ny = 3.30 + row * 1.05
    add_text(s6, nx, ny, col_w_half / 2 - 0.2, 0.6, n,
             font=FONT_MONO, size=30, color=INK, bold=True)
    add_text(s6, nx, ny + 0.65, col_w_half / 2 - 0.2, 0.3, lbl,
             font=FONT_MONO, size=9, color=INK_4)

# Live bullets
bullet_y = 5.6
bullets = [
    "@kyvernlabs/sdk live on npm — 4 lines, 3 framework adapters",
    "ParallaxPay agent integrated — third-party, real x402 origin",
    "Policy program deployed: PpmZErWfT5zpeo1fJtTbpqezFGbRUamaNNRWViaMSqc",
]
for i, b in enumerate(bullets):
    add_text(s6, left_x, bullet_y + i * 0.30, col_w_half, 0.28,
             "✓  " + b, font=FONT_SANS, size=10.5, color=INK_3)

# RIGHT — Roadmap
add_text(s6, right_x, 2.85, col_w_half, 0.3,
         "WHAT'S NEXT",
         font=FONT_MONO, size=10, color=INK_4, bold=True)

roadmap = [
    ("NOW",
     "Per-agent policy layer on devnet · live",
     INK),
    ("NEXT",
     "Mainnet · Kyvern Shield · 100 paid builders",
     INK),
    ("LATER",
     "Kyvern Device — hardware where your agent lives · agent playground "
     "(every agent pays every agent via x402) · closed-loop AI economy.",
     INK),
]
ry = 3.30
for label, body, color in roadmap:
    add_text(s6, right_x, ry, 1.0, 0.4,
             label, font=FONT_MONO, size=11, color=INK, bold=True)
    add_text(s6, right_x + 1.05, ry, col_w_half - 1.1, 0.8,
             body, font=FONT_SANS, size=12.5, color=INK_2, line_spacing=1.4)
    ry += 1.2

# bottom line
add_rule(s6, 0.5, 6.95, SLIDE_W - 1.0)

add_bottom_url_line(s6)
add_text(s6, SLIDE_W - 1.5, 0.32, 1.4, 0.25, "06 / 07",
         font=FONT_MONO, size=8, color=INK_4, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 7 — THE ASK (softened for the room — partnership > number)
# ════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
set_solid_bg(s7, CREAM)
add_top_brand_line(s7)
add_rule(s7, 0.5, 0.76, SLIDE_W - 1.0)

add_text(s7, 0.5, 1.3, 8, 0.4, "SLIDE 07 · THE ASK",
         font=FONT_MONO, size=10, color=INK_4)

# Big headline — what the room can give
add_text(s7, 0.5, 1.95, 12.3, 1.4,
         "Build with us.",
         font=FONT_DISPLAY, size=82, color=INK, bold=True, line_spacing=1.0)

# Three asks — what we actually need from this room
asks = [
    ("BUILDERS",
     "Integrate the SDK. Ship an agent under Kyvern policy."),
    ("PARTNERS",
     "KAST, pay.sh, Solana ecosystem teams — let's talk integrations."),
    ("FEEDBACK",
     "I built this for builders. Tell me what's missing."),
]
ay = 3.65
aw = (SLIDE_W - 1.0) / 3
for i, (head, sub) in enumerate(asks):
    ax = 0.5 + i * aw
    add_text(s7, ax, ay, aw - 0.3, 0.4,
             head, font=FONT_MONO, size=12, color=INK, bold=True)
    add_text(s7, ax, ay + 0.5, aw - 0.3, 0.9,
             sub, font=FONT_SANS, size=13, color=INK_2,
             line_spacing=1.4)

# Quiet line about funding (no number on stage)
add_text(s7, 0.5, 5.1, 12.3, 0.4,
         "Raising a pre-seed privately. Happy to talk after dinner.",
         font=FONT_SANS, size=12, color=INK_3, italic=True)

# Founder line + bio
add_rule(s7, 0.5, 5.65, SLIDE_W - 1.0)

add_text(s7, 0.5, 5.85, 6.5, 0.4,
         "SHARIQ SHAUKAT",
         font=FONT_MONO, size=11, color=INK, bold=True)
add_text(s7, 0.5, 6.20, 8, 0.4,
         "Solo founder · Three x402 products before this. Multiple hackathon wins.",
         font=FONT_SANS, size=12, color=INK_3)

# Contact column right
add_text(s7, 8.0, 5.85, 5, 0.3, "CONTACT",
         font=FONT_MONO, size=9, color=INK_4, bold=True, align=PP_ALIGN.RIGHT)
add_text(s7, 8.0, 6.15, 5, 0.3, "kyvernlabs.com  ·  @shariqshkt",
         font=FONT_MONO, size=11, color=INK, align=PP_ALIGN.RIGHT)
add_text(s7, 8.0, 6.45, 5, 0.3, "shariqshaukat786@gmail.com",
         font=FONT_MONO, size=10, color=INK_3, align=PP_ALIGN.RIGHT)

# Slide marker
add_text(s7, SLIDE_W - 1.5, 0.32, 1.4, 0.25, "07 / 07",
         font=FONT_MONO, size=8, color=INK_4, align=PP_ALIGN.RIGHT)
add_text(s7, 0.5, SLIDE_H - 0.45, 6, 0.25,
         "BUILD ON SOLANA WITH KAST · PAKISTAN · DEMO DAY 2026.06.02",
         font=FONT_MONO, size=8, color=INK_4, align=PP_ALIGN.LEFT)


prs.save(OUT)
print(f"✓  Deck written: {OUT}")
print(f"   {OUT.stat().st_size // 1024} KB · {len(prs.slides)} slides")
